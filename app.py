import os
import sqlite3
import logging
from datetime import datetime
from flask import Flask, render_template, request, redirect, url_for, send_file, g
from io import BytesIO

logging.basicConfig(level=logging.DEBUG)

app = Flask(__name__)
app.secret_key = os.environ.get("SESSION_SECRET", "proposal-secret-key-2024")

DATABASE = "database.db"
ADMIN_SECRET = "julisunkan"


def get_db():
    db = getattr(g, "_database", None)
    if db is None:
        db = g._database = sqlite3.connect(DATABASE)
        db.row_factory = sqlite3.Row
    return db


@app.teardown_appcontext
def close_connection(exception):
    db = getattr(g, "_database", None)
    if db is not None:
        db.close()


def init_db():
    with app.app_context():
        db = sqlite3.connect(DATABASE)
        db.execute("""
            CREATE TABLE IF NOT EXISTS settings (
                id INTEGER PRIMARY KEY,
                api_key TEXT
            )
        """)
        db.execute("""
            CREATE TABLE IF NOT EXISTS proposals (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                business_name TEXT,
                service TEXT,
                target_market TEXT,
                proposal_type TEXT,
                content TEXT,
                created_at TEXT
            )
        """)
        db.execute("""
            CREATE TABLE IF NOT EXISTS reports (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                proposal_id INTEGER,
                reason TEXT,
                reported_at TEXT,
                FOREIGN KEY (proposal_id) REFERENCES proposals(id)
            )
        """)
        # Migrate: add missing columns if upgrading from old schema
        try:
            db.execute("ALTER TABLE proposals ADD COLUMN service TEXT")
        except Exception:
            pass
        try:
            db.execute("ALTER TABLE proposals ADD COLUMN target_market TEXT")
        except Exception:
            pass
        db.commit()
        db.close()


def get_api_key():
    db = get_db()
    row = db.execute("SELECT api_key FROM settings WHERE id = 1").fetchone()
    return row["api_key"] if row else None


def mask_key(key):
    if not key:
        return None
    if len(key) <= 4:
        return "*" * len(key)
    return "*" * (len(key) - 4) + key[-4:]


def generate_proposal_ai(business_name, service, target_market, proposal_type, api_key):
    from groq import Groq
    client = Groq(api_key=api_key)

    prompt = f"""You are a professional business consultant. Write a comprehensive, well-structured business proposal for the following:

Business Name: {business_name}
Service / Product: {service}
Target Market: {target_market}
Proposal Type: {proposal_type}

Generate a full, professional business proposal with these exact sections, using markdown-style headers (## for each section):

## Executive Summary
## Business Overview
## Market Analysis
## Strategy and Implementation
## Budget Estimate
## Expected Results
## Conclusion

Make each section detailed, specific, and professional. Use the provided business details throughout. Write at least 3-4 paragraphs per section."""

    chat_completion = client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama-3.3-70b-versatile",
    )
    return chat_completion.choices[0].message.content


# Serve service worker from root so it has full-site scope
@app.route("/service-worker.js")
def service_worker():
    sw_path = os.path.join(app.root_path, "static", "service-worker.js")
    return send_file(sw_path, mimetype="application/javascript")


@app.route("/", methods=["GET", "POST"])
def index():
    message = None
    msg_type = None

    if request.method == "POST":
        business_name = request.form.get("business_name", "").strip()
        service = request.form.get("service", "").strip()
        target_market = request.form.get("target_market", "").strip()
        proposal_type = request.form.get("proposal_type", "Startup Proposal")

        if not business_name or not service or not target_market:
            message = "Please fill in all fields."
            msg_type = "error"
        else:
            api_key = get_api_key()
            if not api_key:
                message = "No API key configured. Please visit the admin page to add your Groq API key."
                msg_type = "error"
            else:
                try:
                    content = generate_proposal_ai(business_name, service, target_market, proposal_type, api_key)
                    db = get_db()
                    created_at = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    cursor = db.execute(
                        "INSERT INTO proposals (business_name, service, target_market, proposal_type, content, created_at) VALUES (?, ?, ?, ?, ?, ?)",
                        (business_name, service, target_market, proposal_type, content, created_at)
                    )
                    db.commit()
                    proposal_id = cursor.lastrowid
                    return redirect(url_for("edit", proposal_id=proposal_id, success=1))
                except Exception as e:
                    logging.error(f"Error generating proposal: {e}")
                    message = f"Error generating proposal: {str(e)}"
                    msg_type = "error"

    return render_template("index.html", message=message, msg_type=msg_type)


@app.route("/edit/<int:proposal_id>", methods=["GET", "POST"])
def edit(proposal_id):
    db = get_db()
    proposal = db.execute("SELECT * FROM proposals WHERE id = ?", (proposal_id,)).fetchone()

    if not proposal:
        return redirect(url_for("index"))

    message = None
    msg_type = None

    if request.args.get("success"):
        message = "Proposal generated successfully!"
        msg_type = "success"

    if request.method == "POST":
        new_content = request.form.get("content", "")
        db.execute("UPDATE proposals SET content = ? WHERE id = ?", (new_content, proposal_id))
        db.commit()
        message = "Proposal saved successfully!"
        msg_type = "success"
        proposal = db.execute("SELECT * FROM proposals WHERE id = ?", (proposal_id,)).fetchone()

    return render_template("edit.html", proposal=proposal, message=message, msg_type=msg_type)


@app.route("/download-docx/<int:proposal_id>")
def download_docx(proposal_id):
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    db = get_db()
    proposal = db.execute("SELECT * FROM proposals WHERE id = ?", (proposal_id,)).fetchone()
    if not proposal:
        return redirect(url_for("index"))

    doc = Document()

    title = doc.add_heading(f"{proposal['business_name']} - {proposal['proposal_type']}", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph(f"Generated: {proposal['created_at']}")
    if proposal["service"]:
        doc.add_paragraph(f"Service / Product: {proposal['service']}")
    if proposal["target_market"]:
        doc.add_paragraph(f"Target Market: {proposal['target_market']}")
    doc.add_paragraph("")

    for line in proposal["content"].split("\n"):
        line = line.strip()
        if not line:
            doc.add_paragraph("")
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=1)
        elif line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        else:
            doc.add_paragraph(line)

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)

    filename = f"{proposal['business_name'].replace(' ', '_')}_proposal.docx"
    return send_file(buffer, as_attachment=True, download_name=filename,
                     mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")


@app.route("/download-pdf/<int:proposal_id>")
def download_pdf(proposal_id):
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

    db = get_db()
    proposal = db.execute("SELECT * FROM proposals WHERE id = ?", (proposal_id,)).fetchone()
    if not proposal:
        return redirect(url_for("index"))

    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter,
                            rightMargin=inch, leftMargin=inch,
                            topMargin=inch, bottomMargin=inch)

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("CustomTitle", parent=styles["Title"],
                                 fontSize=18, textColor=colors.HexColor("#1e3a5f"), spaceAfter=6)
    heading_style = ParagraphStyle("CustomHeading", parent=styles["Heading1"],
                                   fontSize=13, textColor=colors.HexColor("#16a34a"), spaceBefore=14, spaceAfter=4)
    body_style = ParagraphStyle("CustomBody", parent=styles["Normal"],
                                fontSize=10, leading=15, spaceAfter=6)
    sub_style = ParagraphStyle("CustomSub", parent=styles["Normal"],
                               fontSize=9, textColor=colors.gray, spaceAfter=4)

    story = []
    story.append(Paragraph(f"{proposal['business_name']} — {proposal['proposal_type']}", title_style))
    story.append(Paragraph(f"Generated: {proposal['created_at']}", sub_style))
    if proposal["service"]:
        story.append(Paragraph(f"Service / Product: {proposal['service']}", sub_style))
    if proposal["target_market"]:
        story.append(Paragraph(f"Target Market: {proposal['target_market']}", sub_style))
    story.append(Spacer(1, 0.2 * inch))

    for line in proposal["content"].split("\n"):
        line = line.strip()
        if not line:
            story.append(Spacer(1, 0.08 * inch))
        elif line.startswith("## "):
            story.append(Paragraph(line[3:], heading_style))
        elif line.startswith("# "):
            story.append(Paragraph(line[2:], heading_style))
        else:
            story.append(Paragraph(line, body_style))

    doc.build(story)
    buffer.seek(0)

    filename = f"{proposal['business_name'].replace(' ', '_')}_proposal.pdf"
    return send_file(buffer, as_attachment=True, download_name=filename, mimetype="application/pdf")


@app.route("/download-ppt/<int:proposal_id>")
def download_ppt(proposal_id):
    from pptx import Presentation
    from pptx.util import Pt

    db = get_db()
    proposal = db.execute("SELECT * FROM proposals WHERE id = ?", (proposal_id,)).fetchone()
    if not proposal:
        return redirect(url_for("index"))

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    # Title slide
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = proposal["business_name"]
    subtitle_parts = [proposal["proposal_type"], f"Generated: {proposal['created_at']}"]
    if proposal["service"]:
        subtitle_parts.append(f"Service: {proposal['service']}")
    if proposal["target_market"]:
        subtitle_parts.append(f"Target Market: {proposal['target_market']}")
    slide.placeholders[1].text = "\n".join(subtitle_parts)

    current_section = None
    section_content = []

    def add_slide(title, content_lines):
        s = prs.slides.add_slide(content_slide_layout)
        s.shapes.title.text = title
        tf = s.placeholders[1].text_frame
        tf.clear()
        tf.word_wrap = True
        tf.text = "\n".join(content_lines)
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(11)

    for line in proposal["content"].split("\n"):
        line = line.strip()
        if line.startswith("## ") or line.startswith("# "):
            if current_section and section_content:
                add_slide(current_section, section_content)
            current_section = line.lstrip("# ").strip()
            section_content = []
        elif line:
            section_content.append(line)

    if current_section and section_content:
        add_slide(current_section, section_content)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    filename = f"{proposal['business_name'].replace(' ', '_')}_proposal.pptx"
    return send_file(buffer, as_attachment=True, download_name=filename,
                     mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")


@app.route("/report/<int:proposal_id>", methods=["POST"])
def report(proposal_id):
    db = get_db()
    proposal = db.execute("SELECT * FROM proposals WHERE id = ?", (proposal_id,)).fetchone()
    if not proposal:
        return redirect(url_for("index"))

    reason = request.form.get("reason", "").strip()
    if not reason:
        return redirect(url_for("edit", proposal_id=proposal_id))

    reported_at = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    db.execute(
        "INSERT INTO reports (proposal_id, reason, reported_at) VALUES (?, ?, ?)",
        (proposal_id, reason, reported_at)
    )
    db.commit()
    return redirect(url_for("edit", proposal_id=proposal_id, reported=1))


@app.route("/dashboard")
def dashboard():
    db = get_db()
    proposals = db.execute("SELECT * FROM proposals ORDER BY created_at DESC").fetchall()
    return render_template("dashboard.html", proposals=proposals)


@app.route("/admin/delete-report/<int:report_id>", methods=["POST"])
def delete_report(report_id):
    key = request.args.get("key", "")
    if key != ADMIN_SECRET:
        return redirect(url_for("admin"))
    db = get_db()
    db.execute("DELETE FROM reports WHERE id = ?", (report_id,))
    db.commit()
    return redirect(url_for("admin") + "?key=" + ADMIN_SECRET + "&deleted=1")


@app.route("/admin/delete-proposal/<int:proposal_id>", methods=["POST"])
def delete_proposal(proposal_id):
    key = request.args.get("key", "")
    if key != ADMIN_SECRET:
        return redirect(url_for("admin"))
    db = get_db()
    db.execute("DELETE FROM reports WHERE proposal_id = ?", (proposal_id,))
    db.execute("DELETE FROM proposals WHERE id = ?", (proposal_id,))
    db.commit()
    return redirect(url_for("admin") + "?key=" + ADMIN_SECRET + "&deleted=1")


@app.route("/admin", methods=["GET", "POST"])
def admin():
    key = request.args.get("key", "")
    if key != ADMIN_SECRET:
        return render_template("admin.html", authorized=False, message=None, msg_type=None)

    message = None
    msg_type = None

    if request.args.get("deleted"):
        message = "Deleted successfully."
        msg_type = "success"

    if request.method == "POST":
        api_key = request.form.get("api_key", "").strip()
        if api_key:
            db = get_db()
            existing = db.execute("SELECT id FROM settings WHERE id = 1").fetchone()
            if existing:
                db.execute("UPDATE settings SET api_key = ? WHERE id = 1", (api_key,))
            else:
                db.execute("INSERT INTO settings (id, api_key) VALUES (1, ?)", (api_key,))
            db.commit()
            message = "API key saved successfully!"
            msg_type = "success"
        else:
            message = "Please enter a valid API key."
            msg_type = "error"

    db = get_db()
    current_key = get_api_key()
    reports = db.execute("""
        SELECT r.id as report_id, r.reason, r.reported_at,
               p.id as proposal_id, p.business_name, p.proposal_type, p.created_at
        FROM reports r
        JOIN proposals p ON r.proposal_id = p.id
        ORDER BY r.reported_at DESC
    """).fetchall()
    return render_template("admin.html", authorized=True, message=message,
                           msg_type=msg_type, masked_key=mask_key(current_key),
                           reports=reports)


if __name__ == "__main__":
    init_db()
    app.run(host="0.0.0.0", port=5000, debug=True)
